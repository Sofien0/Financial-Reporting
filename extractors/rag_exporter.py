import os
from pathlib import Path
from datetime import datetime
from textwrap import wrap
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT
from reportlab.lib import colors
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

OUTPUT_DIR = Path("data/outputs")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

def _default_filename(ext: str) -> str:
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return str(OUTPUT_DIR / f"rag_session_{timestamp}.{ext}")

# -------------------
# PDF Export
# -------------------
def export_to_pdf(chat_log, filename=None):
    filename = filename or _default_filename("pdf")
    doc = SimpleDocTemplate(filename, pagesize=A4)
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="Question", fontSize=12, leading=14, spaceAfter=6, textColor=colors.HexColor("#004c99"), leftIndent=0, alignment=TA_LEFT))
    styles.add(ParagraphStyle(name="Answer", fontSize=11, leading=14, spaceAfter=10, leftIndent=12))

    elements = []
    elements.append(Paragraph("ESG RAG Session Export", styles['Title']))
    elements.append(Spacer(1, 20))

    for i, entry in enumerate(chat_log, 1):
        q = f"Q{i}: {entry['question']}"
        elements.append(Paragraph(q, styles['Question']))

        # Bullet-point answers (split on numbers, semicolons, or periods)
        parts = [p.strip() for p in entry['answer'].replace("\n", " ").split(". ") if p.strip()]
        bullets = [ListItem(Paragraph(f"{p}.", styles['Answer'])) for p in parts]
        elements.append(ListFlowable(bullets, bulletType='bullet'))
        elements.append(Spacer(1, 12))

    doc.build(elements)
    return filename

# -------------------
# DOCX Export
# -------------------
def export_to_docx(chat_log, filename=None):
    filename = filename or _default_filename("docx")
    doc = Document()
    doc.add_heading("ESG RAG Session Export", level=0)

    for i, entry in enumerate(chat_log, 1):
        doc.add_heading(f"Q{i}: {entry['question']}", level=1)

        # Bullet-point answers
        parts = [p.strip() for p in entry['answer'].replace("\n", " ").split(". ") if p.strip()]
        for p in parts:
            doc.add_paragraph(p + ".", style="List Bullet")

    doc.save(filename)
    return filename

# -------------------
# PPTX Export
# -------------------
def export_to_pptx(chat_log, filename=None):
    filename = filename or _default_filename("pptx")
    prs = Presentation()

    # --- Title Slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "ESG RAG Session Export"
    slide.placeholders[1].text = datetime.now().strftime("%Y-%m-%d %H:%M")

    # --- Design settings ---
    title_font_size = Pt(28)
    body_font_size = Pt(18)
    max_chars_per_slide = 400

    for i, entry in enumerate(chat_log, 1):
        question = f"Q{i}: {entry['question']}"
        answer = entry['answer']

        # Split into chunks
        chunks = wrap(answer, max_chars_per_slide, break_long_words=False, replace_whitespace=False)

        for j, chunk in enumerate(chunks, 1):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = question if j == 1 else f"{question} (cont.)"

            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()

            # Bullet-point sentences
            sentences = [s.strip() for s in chunk.split(". ") if s.strip()]
            for s in sentences:
                p = tf.add_paragraph()
                p.text = s + "."
                p.font.size = body_font_size

            # Style title
            title_tf = slide.shapes.title.text_frame
            title_tf.paragraphs[0].font.size = title_font_size
            title_tf.paragraphs[0].font.bold = True
            title_tf.paragraphs[0].font.color.rgb = RGBColor(0, 76, 153)  # deep blue

            # Background color
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(240, 248, 255)  # light blue

    prs.save(filename)
    return filename
