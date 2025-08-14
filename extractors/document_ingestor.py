import os
import pandas as pd
from pathlib import Path
from typing import List, Dict, Union
import pdfplumber
import docx
from pptx import Presentation
import nltk
nltk.download('punkt')
from nltk.tokenize import sent_tokenize


def chunk_text(text: str, max_tokens: int = 300) -> List[str]:
    sentences = sent_tokenize(text)
    chunks, current_chunk = [], []
    current_length = 0

    for sentence in sentences:
        current_chunk.append(sentence)
        current_length += len(sentence.split())

        if current_length >= max_tokens:
            chunks.append(" ".join(current_chunk))
            current_chunk = []
            current_length = 0

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks


def load_pdf(path: Path) -> List[Dict]:
    chunks = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text() or ""
                for j, chunk in enumerate(chunk_text(text)):
                    chunks.append({
                        "content": chunk,
                        "source": path.name,
                        "page": i + 1,
                        "chunk_id": f"{path.stem}_p{i+1}_c{j+1}"
                    })
    except Exception as e:
        print(f"❌ Failed to process PDF {path}: {e}")
    return chunks


def load_docx(path: Path) -> List[Dict]:
    chunks = []
    try:
        doc = docx.Document(path)
        full_text = "\n".join(p.text for p in doc.paragraphs)
        for j, chunk in enumerate(chunk_text(full_text)):
            chunks.append({
                "content": chunk,
                "source": path.name,
                "page": None,
                "chunk_id": f"{path.stem}_c{j+1}"
            })
    except Exception as e:
        print(f"❌ Failed to process DOCX {path}: {e}")
    return chunks


def load_pptx(path: Path) -> List[Dict]:
    chunks = []
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            slide_text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            for j, chunk in enumerate(chunk_text(slide_text)):
                chunks.append({
                    "content": chunk,
                    "source": path.name,
                    "page": i + 1,
                    "chunk_id": f"{path.stem}_slide{i+1}_c{j+1}"
                })
    except Exception as e:
        print(f"❌ Failed to process PPTX {path}: {e}")
    return chunks


def load_and_chunk_document(path: Union[str, Path]) -> List[Dict]:
    path = Path(path)
    ext = path.suffix.lower()

    if ext == ".pdf":
        return load_pdf(path)
    elif ext == ".docx":
        return load_docx(path)
    elif ext == ".pptx":
        return load_pptx(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def extract_metadata(path: Path, root: Path) -> Dict:
    relative_path = path.relative_to(root)
    sector = relative_path.parts[0]
    subsector = relative_path.parts[1]
    filename = relative_path.name
    company_name = filename.rsplit("__", 1)[0].replace("_", " ")
    year = filename.rsplit("__", 1)[-1].split(".")[0]
    return {
        "sector": sector,
        "subsector": subsector,
        "company": company_name,
        "year": year,
        "source_path": str(relative_path)
    }


def batch_chunk_all_documents(root: Path, output_path: Path):
    supported_exts = {".pdf", ".docx", ".pptx"}
    all_chunks = []
    total_files = 0

    print(f"🚀 Starting batch ingestion from: {root}")

    for sector_dir in root.iterdir():
        if not sector_dir.is_dir():
            continue
        for subsector_dir in sector_dir.iterdir():
            if not subsector_dir.is_dir():
                continue
            for file in subsector_dir.iterdir():
                if file.suffix.lower() not in supported_exts:
                    continue

                print(f"📄 Processing: {file.relative_to(root)}")
                chunks = load_and_chunk_document(file)
                meta = extract_metadata(file, root)

                for chunk in chunks:
                    chunk.update(meta)
                    all_chunks.append(chunk)

                print(f"   ➕ {len(chunks)} chunks extracted.")
                total_files += 1

    print(f"\n✅ Ingestion complete: {total_files} files, {len(all_chunks)} total chunks.")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    df = pd.DataFrame(all_chunks)
    df.to_parquet(output_path, index=False)
    print(f"💾 Saved chunks to: {output_path}")


if __name__ == "__main__":
    root = Path("data/raw")
    output = Path("data/processed/chunked_documents.parquet")

    batch_chunk_all_documents(root, output)
